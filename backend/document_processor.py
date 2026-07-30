import os
import zipfile
from typing import Dict, List, Optional
from pathlib import Path


def extract_pdf_text(file_path: str) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return "\n\n".join(text_parts)
    except Exception as e:
        return f"[PDF extraction error: {e}]"


def extract_docx_text(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            tables.append("\n".join(rows))
        content = "\n\n".join(paragraphs)
        if tables:
            content += "\n\n--- TABLES ---\n\n" + "\n\n".join(tables)
        return content
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def extract_pptx_text(file_path: str) -> str:
    """Text from a PowerPoint deck, including tables and speaker notes.

    Pre-proposal conference decks are routinely included in solicitation
    packages and sometimes state submission instructions or deadline changes.
    A real 39-file package (MDOT TSOOPD2609) had its pre-proposal slides come
    back as "[Unsupported file type: .pptx]" — silently invisible to the
    checklist. If python-pptx isn't installed the message says so explicitly
    rather than implying the file was read and found empty.
    """
    try:
        from pptx import Presentation
    except ImportError:
        return ("[PPTX not extracted: python-pptx is not installed on this server — "
                "this file's content was NOT reviewed. Convert it to PDF and re-upload.]")
    try:
        prs = Presentation(file_path)
        slides_text = []
        for index, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        row_str = " | ".join(cells).strip(" |")
                        if row_str:
                            parts.append(row_str)
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Speaker notes] {notes}")
            except Exception:
                pass
            if parts:
                slides_text.append(f"[Slide {index}]\n" + "\n".join(parts))
        return "\n\n".join(slides_text)
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def extract_xlsx_text(file_path: str) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                row_str = " | ".join(cells).strip(" |")
                if row_str.strip():
                    rows_text.append(row_str)
            if rows_text:
                sheets_text.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text))
        return "\n\n".join(sheets_text)
    except Exception as e:
        return f"[XLSX extraction error: {e}]"


def extract_zip_contents(file_path: str, extract_dir: str) -> List[Dict]:
    extracted = []
    try:
        # Create the target directory rather than assuming it exists. Render's
        # filesystem is ephemeral, so UPLOAD_PATH can be gone after a restart or
        # redeploy — the first inner file then failed with ENOENT, the whole zip
        # came back as a 237-char "[ZIP extraction error: ...]" string, and that
        # error text was handed to the model AS the solicitation. With nothing
        # real to read it produced generic boilerplate (W-9, capability
        # statement, certificate of insurance) marked on_file, which auto-checks
        # and reports the bid ready to submit. Observed on a real opportunity.
        os.makedirs(extract_dir, exist_ok=True)
        with zipfile.ZipFile(file_path, "r") as zf:
            for name in zf.namelist():
                if name.endswith("/"):
                    continue
                ext = Path(name).suffix.lower()
                if ext not in {".pdf", ".docx", ".doc", ".xlsx", ".xls", ".txt", ".pptx", ".ppt"}:
                    continue
                safe_name = os.path.basename(name)
                out_path = os.path.join(extract_dir, safe_name)
                with zf.open(name) as src, open(out_path, "wb") as dst:
                    dst.write(src.read())
                text = extract_file_text(out_path, ext)
                extracted.append({"filename": safe_name, "text": text})
    except Exception as e:
        extracted.append({"filename": "zip_error", "text": f"[ZIP extraction error: {e}]"})
    return extracted


def extract_file_text(file_path: str, ext: Optional[str] = None) -> str:
    if ext is None:
        ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return extract_pdf_text(file_path)
    elif ext in (".docx", ".doc"):
        return extract_docx_text(file_path)
    elif ext in (".xlsx", ".xls"):
        return extract_xlsx_text(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_pptx_text(file_path)
    elif ext == ".txt":
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        except Exception as e:
            return f"[Text read error: {e}]"
    return f"[Unsupported file type: {ext}]"


def process_document(file_path: str, upload_dir: str, file_content: bytes = None) -> str:
    """Process a document file. If file_path doesn't exist on disk but file_content
    (bytes from DB) is provided, writes to a temp file first — handles cloud deploys
    where the local filesystem is ephemeral."""
    import tempfile
    ext = Path(file_path).suffix.lower()
    if not os.path.exists(file_path) and file_content:
        tmp = tempfile.NamedTemporaryFile(suffix=ext, delete=False)
        try:
            tmp.write(file_content)
            tmp.close()
            return process_document(tmp.name, upload_dir)
        finally:
            try:
                os.unlink(tmp.name)
            except Exception:
                pass
    if ext == ".zip":
        contents = extract_zip_contents(file_path, upload_dir)
        combined = []
        for item in contents:
            combined.append(f"=== {item['filename']} ===\n{item['text']}")
        return "\n\n".join(combined)
    return extract_file_text(file_path, ext)


_EXTRACTION_FAILURE_MARKERS = (
    "[ZIP extraction error", "[PDF extraction error", "[DOCX extraction error",
    "[XLSX extraction error", "[PPTX extraction error", "[Text read error",
    "[Unsupported file type", "[PPTX not extracted",
)


def extraction_failed(text: str) -> bool:
    """True when `text` is an extractor error rather than document content.

    These strings were being passed to the model as though they were the
    solicitation. A model handed 237 characters of error text still returns a
    confident checklist — of invented, generic requirements. Callers must check
    this before treating extracted text as real, so a failed read surfaces as a
    failure instead of as a plausible-looking checklist.
    """
    if not text or not text.strip():
        return True
    head = text.strip()[:400]
    return any(marker in head for marker in _EXTRACTION_FAILURE_MARKERS)


def truncate_for_ai(text: str, max_chars: int = 80000) -> str:
    if len(text) <= max_chars:
        return text
    half = max_chars // 2
    return (
        text[:half]
        + f"\n\n[... {len(text) - max_chars:,} characters truncated ...]\n\n"
        + text[-half:]
    )
